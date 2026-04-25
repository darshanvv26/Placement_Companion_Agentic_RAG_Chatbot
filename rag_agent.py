import os

# Load .env file FIRST before any other imports
try:
    from dotenv import load_dotenv
    load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), ".env"), override=True)
except ImportError:
    pass  # python-dotenv not installed, use system env vars

os.environ["ANONYMIZED_TELEMETRY"] = "False"
import json
import sys
from typing import List, Dict, TypedDict, Annotated, Optional
from dataclasses import dataclass

from langchain_ollama import ChatOllama
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.messages import BaseMessage, HumanMessage, AIMessage, SystemMessage
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langgraph.graph import StateGraph, END
import time
from chromadb.config import Settings

# ─── Load Company Data ────────────────────────────────────────────────────────
CHUNKS_FILE = "chunks_output.json"
try:
    with open(CHUNKS_FILE, "r") as f:
        RAW_DATA = json.load(f)
        CHUNKS_DB = RAW_DATA.get("chunks", {})
        ALL_COMPANIES = sorted(list(CHUNKS_DB.keys()))
except Exception as e:
    print(f"Warning: Could not load companies from {CHUNKS_FILE}: {e}")
    RAW_DATA = {}
    CHUNKS_DB = {}
    ALL_COMPANIES = []

def detect_company(text: str) -> Optional[str]:
    """Detects if a company name from our database is mentioned in the text with fuzzy normalization."""
    import re
    def normalize(s: str) -> str:
        # Remove spaces, dots, dashes, and make lowercase
        return re.sub(r'[\s\.\-_,]', '', s).lower()
        
    text_norm = normalize(text)
    
    # Priority 1: Exact normalized match or containment
    for company in ALL_COMPANIES:
        comp_norm = normalize(company)
        # Check if full company name is in text (normalized)
        if comp_norm in text_norm:
            return company
            
    # Priority 2: Base name normalization (handle "ST Micro" -> "STMicrelectronics")
    for company in ALL_COMPANIES:
        base_name = company.split('_')[0]
        base_norm = normalize(base_name)
        if len(base_norm) > 2 and base_norm in text_norm:
            return company
            
    # Priority 3: Common shorthand/typos
    shorthand = {
        "stmicro": "STMicrelectronics_MTech_2026",
        "st micro": "STMicrelectronics_MTech_2026",
        "stmic": "STMicrelectronics_MTech_2026",
        "novartis": "Novartis_MTech_2026"
    }
    for key, val in shorthand.items():
        if key in text.lower():
            return val

    return None

def is_aggregation_query(text: str) -> bool:
    """Detects if the query is asking for a list or count of companies."""
    text_lower = text.lower()
    keywords = ["how many", "list all", "which companies", "show all", "all companies", "list of companies"]
    return any(kw in text_lower for kw in keywords)

def aggregate_companies(query: str) -> List[str]:
    """Returns a list of companies matching a year or general list."""
    # Check for year in query
    import re
    year_match = re.search(r"202\d", query)
    year = year_match.group(0) if year_match else None
    
    if year:
        return [c for c in ALL_COMPANIES if year in c]
    return ALL_COMPANIES

# ─── Configuration ─────────────────────────────────────────────────────────────
CHROMA_PATH = "./chroma_db"
COLLECTION_NAME = "placement_chunks"

# Choose backend: "ollama", "gemini", or "groq"
LLM_BACKEND = os.getenv("LLM_BACKEND", "gemini") 
MODEL_NAME = os.getenv("GEMINI_MODEL", "gemini-flash-latest")

if LLM_BACKEND == "ollama":
    MODEL_NAME = "llama3"
elif LLM_BACKEND == "groq":
    MODEL_NAME = "llama-3.3-70b-versatile"

# AGENT_MODE: "fast" (single pass) or "agentic" (planner -> executor -> critic)
AGENT_MODE = os.getenv("AGENT_MODE", "fast") 

# ─── State Definition ────────────────────────────────────────────────────────
class AgentState(TypedDict):
    query: str
    original_query: str
    chat_history: List[BaseMessage]
    context_chunks: List[str]
    current_answer: str
    critique: str
    iterations: int
    streamer: any # Optional callback for API streaming

# ─── LLM Setup ─────────────────────────────────────────────────────────────
def get_llm():
    if LLM_BACKEND == "ollama":
        return ChatOllama(model=MODEL_NAME, temperature=0, streaming=True)
    elif LLM_BACKEND == "groq":
        from langchain_openai import ChatOpenAI
        api_key = os.getenv("GROQ_API_KEY")
        return ChatOpenAI(
            model=MODEL_NAME, 
            temperature=0, 
            openai_api_key=api_key, 
            base_url="https://api.groq.com/openai/v1"
        )
    else:
        # For planner/critic nodes: use LangChain (non-streaming REST is usually fine)
        from langchain_google_genai import ChatGoogleGenerativeAI
        api_key = os.getenv("GOOGLE_API_KEY")
        # Force REST transport for langchain-google-genai just in case
        return ChatGoogleGenerativeAI(
            model=MODEL_NAME,
            temperature=0,
            google_api_key=api_key,
            transport="rest"
        )

llm = get_llm()

# ─── Shared Resources ────────────────────────────────────────────────────────
print(f"--- INITIALIZING ({LLM_BACKEND.upper()}:{MODEL_NAME}) ---")
print(f"--- MODE: {AGENT_MODE.upper()} ---")

embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

# Explicitly disable telemetry via chromadb Settings
vectorstore = Chroma(
    persist_directory=CHROMA_PATH,
    collection_name=COLLECTION_NAME,
    embedding_function=embeddings,
    client_settings=Settings(anonymized_telemetry=False)
)

# ─── Nodes ──────────────────────────────────────────────────────────────────

def planner_node(state: AgentState):
    """Rewrites query for better retrieval. Skipped in fast mode."""
    global AGENT_MODE
    if AGENT_MODE == "fast":
        return {"query": state["query"], "iterations": 1}
        
    print("--- PLANNER ---")
    prompt = ChatPromptTemplate.from_messages([
        ("system", "You are a Query Optimizer for a RAG system. Convert the user's request into a standalone search query. "
                  "If the user mentions a company, ensure the company name is prominent in the query. "
                  "Output ONLY the optimized query text."),
        MessagesPlaceholder(variable_name="chat_history"),
        ("human", "{query}")
    ])
    chain = prompt | llm
    start_time = time.time()
    response = chain.invoke({"query": state["query"], "chat_history": state["chat_history"]})
    print(f"--- PLANNER FINISHED in {time.time() - start_time:.2f}s ---")
    
    return {
        "query": response.content.strip(),
        "iterations": state.get("iterations", 0) + 1
    }

def retriever_node(state: AgentState):
    """Fetches relevant chunks from ChromaDB or performs aggregation."""
    print("--- RETRIEVER ---")
    start_time = time.time()
    
    query = state["query"]
    original_query = state["original_query"]
    
    # 1. Check for Aggregation Query
    if is_aggregation_query(original_query) or is_aggregation_query(query):
        print("  → Aggregation query detected!")
        companies = aggregate_companies(original_query)
        context = [f"The following companies are in the records: {', '.join(companies)}"]
        return {"context_chunks": context}

    # 2. Targeted Company Search
    company_filter = detect_company(original_query) or detect_company(query)
    
    if company_filter:
        print(f"  → Targeted search for company: {company_filter}")
        # Use local chunks strictly for highest accuracy if company detected
        local_chunks = CHUNKS_DB.get(company_filter, [])
        if local_chunks:
            context = [f"Source: {chunk.get('company', company_filter)}\nContent: {chunk['content']}" for chunk in local_chunks]
            print(f"--- RETRIEVER (LOCAL) FINISHED in {time.time() - start_time:.2f}s ---")
            return {"context_chunks": context}
    
    # 3. Fallback to Semantic Search
    print("  → Semantic search via ChromaDB")
    filter_dict = {"company": company_filter} if company_filter else None
    results = vectorstore.similarity_search(query, k=6 if company_filter else 4, filter=filter_dict)
    
    if not results and filter_dict:
        print("    → No results with filter, trying general search...")
        results = vectorstore.similarity_search(query, k=4)

    context = []
    # Limit to top 6 chunks to prevent context bloat and high latency on CPU
    for doc in results[:6]:
        company_name = doc.metadata.get("company", "Unknown")
        context.append(f"Source: {company_name}\nContent: {doc.page_content}")
        
    print(f"--- RETRIEVER (CHROMA) FINISHED in {time.time() - start_time:.2f}s ---")
    return {"context_chunks": context}

def executor_node(state: AgentState):
    """Generates the main answer using retrieved context with manual REST calls."""
    print("--- EXECUTOR (Manual REST) ---")
    print("─"*30)

    context_text = "\n\n---\n\n".join(state["context_chunks"])
    sliced_history = state["chat_history"][-4:]
    streamer = state.get("streamer")
    start_time = time.time()
    response = ""

    if LLM_BACKEND == "gemini":
        import urllib.request
        import json
        api_key = os.getenv("GOOGLE_API_KEY")
        
        system_instr = (
            "You are the MSIS Placement Assistant. Your goal is to be helpful, polite, and professional while providing accurate information about campus placements.\n\n"
            "### TONE & STYLE RULES:\n"
            "- Start with a very brief polite greeting if it's a new topic (e.g., 'Hello! I'd be happy to help you with that...').\n"
            "- If the user asks for a specific company or detail, respond politely and nicely.\n"
            "- Use clear, scannable Markdown formatting.\n"
            "- Use POINT-WISE LISTS for all details like eligibility, selection rounds, and stipends.\n"
            "- Use BOLD headers for different sections.\n"
            "- Keep a professional and encouraging tone throughout.\n\n"
            "### CRITICAL ACCURACY RULES:\n"
            "- Answer ONLY using the provided CONTEXT. Do NOT use external knowledge.\n"
            "- If the answer is not in context, say politely: 'I'm sorry, but I don't have that information in the current placement records. Please check with the placement office for more details.'\n"
            "- Be extremely precise with numbers (salaries, dates). Only quote exactly what is in the context.\n"
            "- ALWAYS add a '#### 📚 Sources' section at the very end."
        )

        # Build prompt parts for manual REST call
        history_history = []
        for msg in sliced_history:
            role = "user" if msg.__class__.__name__ == "HumanMessage" else "model"
            history_history.append({"role": role, "parts": [{"text": msg.content}]})

        full_prompt = f"### CONTEXT:\n{context_text}\n\n### QUESTION:\n{state['original_query']}"
        
        # Payload for streamGenerateContent
        payload = {
            "contents": history_history + [{"role": "user", "parts": [{"text": full_prompt}]}],
            "systemInstruction": {"parts": [{"text": system_instr}]},
            "generationConfig": {"temperature": 0}
        }
        
        # Native model name usually doesn't need 'models/' in the URL path if using v1beta
        # URL for streaming REST API
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{MODEL_NAME}:streamGenerateContent?key={api_key}"
        
        headers = {
            'Content-Type': 'application/json', 
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko)'
        }
        
        # Try with retry logic for 429
        max_retries = 3
        for attempt in range(max_retries):
            try:
                req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'), headers=headers)
                with urllib.request.urlopen(req, timeout=30) as f:
                    import re
                    for line_bytes in f:
                        line = line_bytes.decode('utf-8')
                        matches = re.findall(r'"text":\s*"((?:[^"\\]|\\.)*)"', line)
                        for match in matches:
                            try:
                                text = json.loads('"' + match + '"')
                                if text:
                                    response += text
                                    if streamer: streamer(text)
                                    else:
                                        sys.stdout.write(text); sys.stdout.flush()
                            except: continue
                break # Success
            except Exception as e:
                if "429" in str(e) and attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 2
                    print(f"Rate limited (429). Retrying in {wait_time}s...")
                    time.sleep(wait_time)
                    continue
                
                if "429" in str(e):
                    error_msg = ("⚠️ **Rate Limit Reached**: The Gemini free-tier quota (15 requests/min) has been hit. "
                                 "Please wait 60 seconds. You can also try a different API key!")
                else:
                    error_msg = f"Error in Gemini REST call: {str(e)}"
                
                print(f"FAILED: {error_msg}")
                response = error_msg
                if streamer: streamer(error_msg)
                break
                
    elif LLM_BACKEND == "groq":
        # ── Groq REST manual implementation ────────────────────────────────
        import urllib.request
        import json
        api_key = os.getenv("GROQ_API_KEY")
        
        system_instr = (
            "You are the MSIS Placement Assistant. Be helpful, professional, and grounded in CONTEXT ONLY.\n"
            "Use point-wise lists and bold headers. Always add Sources at the end."
        )

        messages = [{"role": "system", "content": system_instr}]
        for msg in sliced_history:
            role = "user" if msg.__class__.__name__ == "HumanMessage" else "assistant"
            messages.append({"role": role, "content": msg.content})
        
        messages.append({"role": "user", "content": f"### CONTEXT:\n{context_text}\n\n### QUESTION:\n{state['original_query']}"})
        
        payload = {
            "model": MODEL_NAME,
            "messages": messages,
            "temperature": 0,
            "stream": True # Standard OpenAI-style streaming
        }
        
        url = "https://api.groq.com/openai/v1/chat/completions"
        
        headers = {
            'Content-Type': 'application/json', 
            'Authorization': f'Bearer {api_key}',
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko)'
        }
        
        try:
            req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'), headers=headers)
            
            with urllib.request.urlopen(req, timeout=30) as f:
                for line_bytes in f:
                    line = line_bytes.decode('utf-8').strip()
                    if not line.startswith("data: ") or line == "data: [DONE]": continue
                    
                    try:
                        chunk_json = json.loads(line[6:])
                        token = chunk_json.get("choices", [{}])[0].get("delta", {}).get("content", "")
                        if token:
                            response += token
                            if streamer: streamer(token)
                            else:
                                sys.stdout.write(token); sys.stdout.flush()
                    except: continue

        except Exception as e:
            error_msg = f"Error in Groq call: {str(e)}"
            print(f"FAILED: {error_msg}")
            response = error_msg
            if streamer: streamer(error_msg)
            
    else:
        # ── Ollama path ─────────────────────────────────────────────────────
        prompt = ChatPromptTemplate.from_messages([
            ("system",
             "You are the MSIS Placement Assistant. Be helpful, polite, and professional.\n\n"
             "### TONE & FORMATTING:\n"
             "- Use a polite, welcoming tone.\n"
             "- Use POINT-WISE lists for details.\n"
             "- Use Markdown headers for organization.\n"
             "- Start with a brief polite intro.\n\n"
             "### RULES:\n"
             "- Answer ONLY using context below.\n"
             "- If unknown, say: 'I'm sorry, I don't have this in the records.'\n"
             "- ALWAYS add a 'Sources' section.\n\n"
             "CONTEXT:\n{context}"),
            MessagesPlaceholder(variable_name="chat_history"),
            ("human", "{original_query}")
        ])
        chain = prompt | llm
        for chunk in chain.stream({"context": context_text, "chat_history": sliced_history, "original_query": state["original_query"]}):
            content = chunk.content if hasattr(chunk, "content") else str(chunk)
            response += content
            if streamer: streamer(content)
            else:
                sys.stdout.write(content); sys.stdout.flush()

    print("\n" + "─"*30)
    print(f"--- EXECUTOR FINISHED in {time.time() - start_time:.2f}s ---")
    return {"current_answer": response}
    

def analyze_document(file_bytes: bytes, file_type: str, doc_type: str = "resume", streamer=None):
    """
    Analyzes a document using Gemini or Groq multimodal/text capabilities.
    """
    print(f"--- DOCUMENT ANALYZER (Doc: {doc_type}, Backend: {LLM_BACKEND}) ---")
    start_time = time.time()
    
    # 1. Specialized instructions
    if doc_type == "jd":
        system_instr = (
            "You are a Technical Recruiter. Analyze the provided Company Job Description (JD).\n"
            "Summarize the role, list top 5 technical skills, provide a 1-week preparation roadmap, and 5 interview questions."
        )
    else:
        system_instr = (
            "You are a Senior Career Coach. Analyze the provided resume. "
            "Suggest skills to add and match against placement context. Rate out of 10."
        )

    # 2. Handle Groq Analysis (Extraction fallback for PDF)
    if LLM_BACKEND == "groq":
        import urllib.request
        import json
        import base64
        api_key = os.getenv("GROQ_API_KEY")
        
        extracted_text = ""
        if 'application/pdf' in file_type:
            try:
                import io
                import PyPDF2
                reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
                for page in reader.pages:
                    extracted_text += page.extract_text() + "\n"
                print(f"  → PDF Extracted: {len(extracted_text)} chars")
            except Exception as pe:
                extracted_text = f"[Error extracting PDF: {str(pe)}]"
        elif 'wordprocessingml' in file_type:
            try:
                import io
                import docx
                doc = docx.Document(io.BytesIO(file_bytes))
                extracted_text = "\n".join([p.text for p in doc.paragraphs])
                print(f"  → Word (DOCX) Extracted: {len(extracted_text)} chars")
            except Exception as de:
                extracted_text = f"[Error extracting DOCX: {str(de)}]"
        elif 'presentationml' in file_type:
            try:
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
                print(f"  → PPTX Extracted: {len(extracted_text)} chars")
            except Exception as pte:
                extracted_text = f"[Error extracting PPTX: {str(pte)}]"
        
        # If image, we could theoretically use Llama 3.2 Vision on Groq
        # For now, we'll try and send as image if it's an image
        if 'image' in file_type:
            # Groq Llama 3.2 Vision supports base64 images in OpenAI format
            encoded_image = base64.b64encode(file_bytes).decode('utf-8')
            messages = [
                {"role": "system", "content": system_instr},
                {"role": "user", "content": [
                    {"type": "text", "text": "Please analyze this document image."},
                    {"type": "image_url", "image_url": {"url": f"data:{file_type};base64,{encoded_image}"}}
                ]}
            ]
            model_to_use = "llama-3.2-11b-vision-preview"
        else:
            messages = [
                {"role": "system", "content": system_instr},
                {"role": "user", "content": f"Analyze this text from a {doc_type} document:\n\n{extracted_text}"}
            ]
            model_to_use = MODEL_NAME

        payload = {
            "model": model_to_use,
            "messages": messages,
            "temperature": 0.2,
            "stream": True
        }
        
        response = ""
        url = "https://api.groq.com/openai/v1/chat/completions"
        headers = {
            'Content-Type': 'application/json', 
            'Authorization': f'Bearer {api_key}',
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        try:
            req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'), headers=headers)
            with urllib.request.urlopen(req, timeout=60) as f:
                for line_bytes in f:
                    line = line_bytes.decode('utf-8').strip()
                    if not line.startswith("data: ") or line == "data: [DONE]": continue
                    try:
                        chunk_json = json.loads(line[6:])
                        token = chunk_json.get("choices", [{}])[0].get("delta", {}).get("content", "")
                        if token:
                            response += token
                            if streamer: streamer(token)
                            else:
                                sys.stdout.write(token); sys.stdout.flush()
                    except: continue
        except Exception as e:
            error_msg = f"Groq analysis error: {str(e)}"
            print(error_msg); response = error_msg
            if streamer: streamer(error_msg)
        return response

    # 3. Handle Gemini Analysis (Multimodal built-in)
    else:
        import urllib.request
        import json
        import base64
        api_key = os.getenv("GOOGLE_API_KEY")

        if doc_type == "resume":
            # Fetch context for resume match
            summary_context = []
            for company, chunks in list(CHUNKS_DB.items())[:15]:
                reqs = [c['content'] for c in chunks if 'eligibility' in c['content'].lower() or 'selection' in c['content'].lower()]
                if reqs:
                    summary_context.append(f"Company: {company}\nData: {' '.join(reqs[:2])}")
            context_text = "\n\n---\n\n".join(summary_context)
            user_prompt = (
                "Here is my resume. Please analyze it based on the following placement records:\n\n"
                f"### COMPANY CONTEXT:\n{context_text}\n\n"
                "Provide a summary, skill recommendations, company matches, and a rating."
            )
        else:
            user_prompt = "Here is a Job Description for a company. Please analyze it and help me prepare for it."

        encoded_file = base64.b64encode(file_bytes).decode('utf-8')
        payload = {
            "contents": [{
                "parts": [
                    {"inline_data": {"mime_type": file_type, "data": encoded_file}},
                    {"text": user_prompt}
                ]
            }],
            "systemInstruction": {"parts": [{"text": system_instr}]},
            "generationConfig": {"temperature": 0.2}
        }

        headers = {
            'Content-Type': 'application/json',
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }

        url = f"https://generativelanguage.googleapis.com/v1beta/models/{MODEL_NAME}:streamGenerateContent?key={api_key}"
        
        response = ""
        try:
            max_retries = 2
            for attempt in range(max_retries):
                try:
                    req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'), headers=headers)
                    with urllib.request.urlopen(req, timeout=60) as f:
                        import re
                        for line_bytes in f:
                            line = line_bytes.decode('utf-8')
                            matches = re.findall(r'"text":\s*"((?:[^"\\]|\\.)*)"', line)
                            for match in matches:
                                try:
                                    text = json.loads('"' + match + '"')
                                    if text:
                                        response += text
                                        if streamer: streamer(text)
                                        else: sys.stdout.write(text); sys.stdout.flush()
                                except: continue
                    break
                except Exception as e:
                    if "429" in str(e) and attempt < max_retries - 1:
                        time.sleep(3)
                        continue
                    raise e
        except Exception as e:
            if "429" in str(e):
                error_msg = "⚠️ **Rate Limit Reached**: The Gemini free-tier quota is currently full. Please wait 60 seconds."
            else:
                error_msg = f"Error analyzing {doc_type}: {str(e)}"
            print(f"FAILED: {error_msg}")
            response = error_msg
            if streamer: streamer(error_msg)
            
        print(f"\n--- {doc_type.upper()} ANALYSIS FINISHED in {time.time() - start_time:.2f}s ---")
        return response

def critic_node(state: AgentState):
    """Validates the answer. Skipped in fast mode."""
    global AGENT_MODE
    if AGENT_MODE == "fast":
        return {"critique": "CORRECT"}

    print("--- CRITIC ---")
    prompt = ChatPromptTemplate.from_messages([
        ("system", "Compare AI Answer against Context. Output 'CORRECT' if perfect. Otherwise, describe issues.\n\nContext:\n{context}\n\nAnswer:\n{answer}"),
        ("human", "Is this answer grounded?")
    ])
    chain = prompt | llm
    context_text = "\n\n---\n\n".join(state["context_chunks"])
    start_time = time.time()
    response = chain.invoke({
        "context": context_text,
        "answer": state["current_answer"]
    })
    print(f"--- CRITIC FINISHED in {time.time() - start_time:.2f}s ---")
    
    return {"critique": response.content}

def router_function(state: AgentState):
    """Decides whether to loop or end."""
    global AGENT_MODE
    if AGENT_MODE == "fast" or "CORRECT" in state["critique"].upper() or state["iterations"] >= 3:
        return "end"
    print(f"Critique found issues: {state['critique'][:50]}... Re-planning.")
    return "replan"

# ─── Graph Construction ─────────────────────────────────────────────────────

workflow = StateGraph(AgentState)

workflow.add_node("planner", planner_node)
workflow.add_node("retriever", retriever_node)
workflow.add_node("executor", executor_node)
workflow.add_node("critic", critic_node)

workflow.set_entry_point("planner")
workflow.add_edge("planner", "retriever")
workflow.add_edge("retriever", "executor")
workflow.add_edge("executor", "critic")

workflow.add_conditional_edges(
    "critic",
    router_function,
    {
        "end": END,
        "replan": "planner"
    }
)

app = workflow.compile()

# ─── Interaction Loop ────────────────────────────────────────────────────────

def main():
    global AGENT_MODE
    print("\n" + "="*50)
    print("Placement Agent (Streaming Enabled)")
    print("="*50)
    print("Commands: /fast, /agentic, /status, exit")
    
    history = []
    
    while True:
        try:
            user_input = input(f"\n[{AGENT_MODE.upper()}] User> ")
        except EOFError:
            break
            
        if user_input.lower() in ["exit", "quit"]:
            break
            
        if user_input.startswith("/"):
            cmd = user_input.lower().strip()
            if cmd == "/fast":
                AGENT_MODE = "fast"
                print("--- Mode switched to FAST ---")
                continue
            elif cmd == "/agentic":
                AGENT_MODE = "agentic"
                print("--- Mode switched to AGENTIC ---")
                continue
            elif cmd == "/status":
                print(f"Backend: {LLM_BACKEND} | Model: {MODEL_NAME} | Mode: {AGENT_MODE}")
                continue
            else:
                print("Unknown command. Use /fast, /agentic, or /status.")
                continue

        initial_state = {
            "query": user_input,
            "original_query": user_input,
            "chat_history": history,
            "context_chunks": [],
            "current_answer": "",
            "critique": "",
            "iterations": 0
        }
        
        # Invoke the graph (streaming happens inside the executor node)
        result = app.invoke(initial_state)
        
        # Update memory
        history.append(HumanMessage(content=user_input))
        history.append(AIMessage(content=result["current_answer"]))
        if len(history) > 10:
            history = history[-10:]

if __name__ == "__main__":
    main()
