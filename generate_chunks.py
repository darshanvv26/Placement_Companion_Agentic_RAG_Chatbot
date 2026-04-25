"""
generate_chunks.py
------------------
Reads all company folders in Placements_Data/ and generates chunks_output.json.
Uses pre-processed OCR results from OCR_mmd_collected/ for image data.

Supported file types:
  .txt   → raw text
  .docx  → python-docx
  .pdf   → pdfplumber
  .pptx  → python-pptx
  .mmd   → Markdown (from OCR_mmd_collected)

Usage:
  python3 generate_chunks.py
"""

import os
import json
import re
from datetime import datetime
from pathlib import Path

# ─── Third-party imports ───────────────────────────────────────────────────────
import docx          # python-docx
import pdfplumber
from pptx import Presentation

# ─── Configuration ─────────────────────────────────────────────────────────────
PLACEMENTS_DIR  = Path(__file__).parent / "Placements_Data"
OCR_MMD_DIR     = Path(__file__).parent / "OCR_mmd_Collected"
OUTPUT_FILE     = Path(__file__).parent / "chunks_output.json"
MAX_CHUNK_CHARS = 1500   # optimized for RAG
CHUNK_OVERLAP   = 150    # character overlap for context preservation

# ─── Section classification keywords ───────────────────────────────────────────
SECTION_KEYWORDS = [
    ("Eligibility Criteria",  ["eligib", "cgpa", " % ", "criteria", "backlog", "gpa", "aggregate", "percentage"]),
    ("Key Responsibilities",  ["key responsib", "your responsib", "key duties"]),
    ("Job Description",       ["job description", "position overview", "role overview", "internship",
                               "location:", "duration:", "stipend:", "job title", "job profile"]),
    ("Skill Set Requirements",["skill", "qualificat", "requirement", "proficien", "tool", "technolog",
                               "experience with", "knowledge of"]),
    ("Compensation",          ["ctc", "compensation", "salary", "package", "lpa", "per annum", "fixed pay"]),
    ("Selection Rounds",      ["round", "interview", "online test", "assessment", "selection process",
                               "test format", "coding round"]),
    ("Benefits",              ["benefit", "insurance", "allowance", "leave", "perk"]),
    ("Program Details",       ["program", "curriculum", "schedule", "semester"]),
    ("Drive Details",         ["drive", "apply before", "deadline", "portal", "spot offer",
                               "starts on", "apply by"]),
]

SKIP_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg",
                   ".ipynb", ".csv", ".html", ".htm", ".rmd", ".nb"}


# ─── 1. Text extraction helpers ────────────────────────────────────────────────

def extract_txt(path: Path) -> str:
    """Read plain text file or .mmd file."""
    try:
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    except Exception as e:
        print(f"  [WARN] text read failed: {path.name} – {e}")
        return ""


def extract_docx(path: Path) -> str:
    """Extract text from .docx, preserving basic structure in markdown."""
    try:
        doc = docx.Document(str(path))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append("")
                continue
            # Headings → markdown heading
            style = para.style.name.lower()
            if "heading 1" in style:
                lines.append(f"# {text}")
            elif "heading 2" in style:
                lines.append(f"## {text}")
            elif "heading 3" in style:
                lines.append(f"### {text}")
            elif "list" in style:
                lines.append(f"- {text}")
            else:
                # Inline bold runs
                md_line = _runs_to_markdown(para)
                lines.append(md_line)
        return "\n".join(lines).strip()
    except Exception as e:
        print(f"  [WARN] docx read failed: {path.name} – {e}")
        return ""


def _runs_to_markdown(para) -> str:
    """Convert a paragraph's runs to markdown-inline bold."""
    parts = []
    for run in para.runs:
        t = run.text
        if not t:
            continue
        if run.bold:
            t = f"**{t.strip()}**"
        parts.append(t)
    return "".join(parts).strip()


def extract_pdf(path: Path) -> str:
    """Extract text from PDF using pdfplumber."""
    try:
        pages = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text.strip())
        return "\n\n".join(pages).strip()
    except Exception as e:
        print(f"  [WARN] pdf read failed: {path.name} – {e}")
        return ""


def extract_pptx(path: Path) -> str:
    """Extract text from .pptx, slide titles as ##, body as bullets."""
    try:
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for j, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    # First paragraph of first shape → treat as slide title
                    if j == 0 and shape == slide.shapes[0]:
                        slide_lines.append(f"## {text}")
                    else:
                        slide_lines.append(f"- {text}")
            if slide_lines:
                slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text).strip()
    except Exception as e:
        print(f"  [WARN] pptx read failed: {path.name} – {e}")
        return ""


def extract_text(path: Path) -> tuple[str, str]:
    """
    Dispatch to the right extractor based on file extension.
    Returns (text, file_type_string).
    """
    ext = path.suffix.lower()
    if ext in (".txt", ".mmd"):
        return extract_txt(path), ext
    elif ext == ".docx":
        return extract_docx(path), ".docx"
    elif ext == ".pdf":
        return extract_pdf(path), ".pdf"
    elif ext == ".pptx":
        return extract_pptx(path), ".pptx"
    else:
        return "", ext


# ─── 2. Markdown formatting pass ───────────────────────────────────────────────

def markdown_format(text: str) -> str:
    """
    Light post-processing to improve markdown readability:
    - ALL-CAPS LABEL: → **ALL-CAPS LABEL:**
    - Lone short lines (likely headings) get ##
    """
    lines = text.split("\n")
    result = []
    for line in lines:
        stripped = line.strip()

        # Already has markdown prefix — skip
        if stripped.startswith(("#", "-", "*", ">", "|")):
            result.append(line)
            continue

        # ALL-CAPS label like "ELIGIBILITY:" or "CTC:" → bold
        if re.match(r'^[A-Z][A-Z &/\-]{2,}:\s', stripped):
            stripped = "**" + stripped.replace(":", ":**", 1)
            result.append(stripped)
            continue

        # Short standalone line with no period → likely a heading
        if 3 < len(stripped) < 60 and not stripped.endswith(".") and "  " not in stripped:
            if stripped.isupper() or stripped.istitle():
                result.append(f"## {stripped}")
                continue

        result.append(line)
    return "\n".join(result)


# ─── 3. Chunking ───────────────────────────────────────────────────────────────

def split_into_chunks(text: str, max_chars: int = MAX_CHUNK_CHARS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Split text into chunks with overlap for better context preservation.
    """
    if not text.strip():
        return []

    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + max_chars
        
        if end < text_len:
            # Try to find a logical break point (double newline, then single newline)
            chunk_slice = text[start : end + 100] # look ahead slightly
            
            # 1. Double newline
            break_point = chunk_slice.rfind("\n\n", 0, max_chars + 1)
            # 2. Single newline
            if break_point == -1:
                break_point = chunk_slice.rfind("\n", 0, max_chars + 1)
            # 3. Last space
            if break_point == -1:
                break_point = chunk_slice.rfind(" ", 0, max_chars + 1)
                
            if break_point != -1 and break_point > max_chars * 0.5:
                end = start + break_point
        
        chunks.append(text[start:end].strip())
        
        # Advance with overlap
        start = end - overlap
        if start < 0: start = 0
        if end >= text_len: break
        
    return [c for c in chunks if c.strip()]


# ─── 4. Section classification ─────────────────────────────────────────────────

def classify_section(text: str) -> str:
    """Classify a chunk's section by keyword matching on lowercase text."""
    lower = text.lower()
    for section_name, keywords in SECTION_KEYWORDS:
        for kw in keywords:
            if kw in lower:
                return section_name
    return "Other"


# ─── 5. Main chunking pipeline ─────────────────────────────────────────────────

def build_header(company: str, section: str) -> str:
    return f"{company} – {section}"


def process_file(
    file_path: Path,
    company: str,
    role: str,
    chunk_id_counter: list,   # mutable counter [int]
) -> list[dict]:
    """Extract, format, split, classify one file into chunks."""
    ext = file_path.suffix.lower()
    if ext in SKIP_EXTENSIONS:
        return []

    print(f"  Processing: {file_path.name}")
    raw_text, file_type = extract_text(file_path)

    if not raw_text.strip():
        print(f"    → Empty / no text extracted")
        return []

    formatted = markdown_format(raw_text)
    text_chunks = split_into_chunks(formatted)

    chunks = []
    for chunk_text in text_chunks:
        if not chunk_text.strip():
            continue
        section = classify_section(chunk_text)
        # result.mmd (OCR) chunks get 0.9 confidence as they are pre-processed
        confidence = 0.9 if file_type == ".mmd" else 1.0

        chunk = {
            "chunk_id":  chunk_id_counter[0],
            "header":    build_header(company, section),
            "section":   section,
            "confidence": confidence,
            "content":   chunk_text,
            "company":   company,
            "role":      role,
            "filename":  file_path.name,
            "file_type": file_type,
        }
        chunks.append(chunk)
        chunk_id_counter[0] += 1

    print(f"    → {len(chunks)} chunk(s)")
    return chunks


def process_company(company_dir: Path) -> list[dict]:
    """Process all files (and sub-role folders) for one company."""
    company = company_dir.name
    print(f"\n[Company] {company}")
    all_chunks = []
    chunk_id_counter = [1]  # per-company counter

    # 1. Process files in company root
    for entry in sorted(company_dir.iterdir()):
        if entry.is_file():
            chunks = process_file(entry, company, "General", chunk_id_counter)
            all_chunks.extend(chunks)

    # 2. Check for local OCR result in company root
    ocr_mmd = OCR_MMD_DIR / company / "result.mmd"
    if ocr_mmd.exists():
        print(f"  Found local OCR: {ocr_mmd.parent.name}/result.mmd")
        chunks = process_file(ocr_mmd, company, "General", chunk_id_counter)
        all_chunks.extend(chunks)

    # 3. Process sub-role folders
    for entry in sorted(company_dir.iterdir()):
        if entry.is_dir():
            role = entry.name
            print(f"  [Role] {role}")
            for file_path in sorted(entry.iterdir()):
                if file_path.is_file():
                    chunks = process_file(file_path, company, role, chunk_id_counter)
                    all_chunks.extend(chunks)
            
            # Check for local OCR result in role folder
            role_ocr_mmd = OCR_MMD_DIR / company / role / "result.mmd"
            if role_ocr_mmd.exists():
                print(f"    Found local OCR for role: {role}/result.mmd")
                chunks = process_file(role_ocr_mmd, company, role, chunk_id_counter)
                all_chunks.extend(chunks)

    print(f"  → Total chunks for {company}: {len(all_chunks)}")
    return all_chunks


# ─── 6. Entry point ────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("Placement Chunk Generator (Local OCR Mode)")
    print("=" * 60)

    if not PLACEMENTS_DIR.exists():
        print(f"ERROR: {PLACEMENTS_DIR} not found")
        return

    output = {
        "generated_at":    datetime.now().isoformat(),
        "total_companies": 0,
        "total_chunks":    0,
        "chunks":          {},
    }

    company_dirs = sorted(
        d for d in PLACEMENTS_DIR.iterdir() if d.is_dir()
    )

    for company_dir in company_dirs:
        chunks = process_company(company_dir)
        if chunks:
            output["chunks"][company_dir.name] = chunks
            output["total_companies"] += 1
            output["total_chunks"] += len(chunks)

    # Write output
    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=2, ensure_ascii=False)

    print("\n" + "=" * 60)
    print(f"Done! Total companies : {output['total_companies']}")
    print(f"      Total chunks     : {output['total_chunks']}")
    print(f"      Output file      : {OUTPUT_FILE}")
    print("=" * 60)


if __name__ == "__main__":
    main()
